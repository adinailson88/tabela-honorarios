# -*- coding: utf-8 -*-
"""Gera APRESENTACAO_CREA_SENGE.pptx (16:9) com python-pptx. Sem dependencias externas."""
import os
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE=os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
DEEP=RGBColor(0x06,0x5A,0x82); TEAL=RGBColor(0x1C,0x72,0x93); MID=RGBColor(0x21,0x29,0x5C)
INK=RGBColor(0x22,0x2A,0x33); MUT=RGBColor(0x5B,0x6B,0x7A); WHITE=RGBColor(0xFF,0xFF,0xFF)
LIGHT=RGBColor(0xF2,0xF5,0xF8); ACC=RGBColor(0x02,0xC3,0x9A)
HSANS='Calibri'; HSERIF='Cambria'

prs=Presentation(); prs.slide_width=I(13.333); prs.slide_height=I(7.5)
BLANK=prs.slide_layouts[6]
SW,SH=13.333,7.5

def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s

def box(s,x,y,w,h,text,size=16,color=INK,bold=False,font=HSANS,align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,italic=False,line=1.05,space=4):
    tb=s.shapes.add_textbox(I(x),I(y),I(w),I(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
    lines=text.split('\n')
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=line; p.space_after=Pt(space)
        r=p.add_run(); r.text=ln; f=r.font
        f.size=Pt(size); f.bold=bold; f.italic=italic; f.name=font; f.color.rgb=color
    return tb

def shape(s,kind,x,y,w,h,fill=None,linec=None,linew=1.0):
    sp=s.shapes.add_shape(kind,I(x),I(y),I(w),I(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if linec is None: sp.line.fill.background()
    else: sp.line.color.rgb=linec; sp.line.width=Pt(linew)
    return sp

def circle_icon(s,x,y,d,fill,glyph,gsize=18):
    shape(s,MSO_SHAPE.OVAL,x,y,d,d,fill=fill)
    box(s,x,y,d,d,glyph,size=gsize,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def stat(s,x,y,w,num,lbl,numc=DEEP):
    box(s,x,y,w,0.9,num,size=40,color=numc,bold=True,font=HSERIF)
    box(s,x,y+0.85,w,0.8,lbl,size=12.5,color=MUT)

# 1 TITLE
s=slide(MID)
shape(s,MSO_SHAPE.OVAL,11.0,-1.4,3.6,3.6,fill=DEEP)
shape(s,MSO_SHAPE.OVAL,12.0,5.3,3.2,3.2,fill=TEAL)
box(s,0.9,2.2,11,1.0,"Nova Metodologia para a Tabela de Honorários",size=40,color=WHITE,bold=True,font=HSERIF)
box(s,0.9,3.25,11,0.9,"dos Profissionais de Engenharia da Bahia",size=30,color=RGBColor(0xCA,0xDC,0xFC),bold=True,font=HSERIF)
box(s,0.92,4.5,11,0.6,"Subsídio técnico baseado em evidências — SENGE/BA",size=17,color=RGBColor(0xCA,0xDC,0xFC))
box(s,0.92,6.4,11,0.5,"Adinailson Guimarães de Oliveira · Engenheiro Eletricista · Conselheiro CREA-BA",size=13,color=RGBColor(0x9F,0xB0,0xC3))

# 2 PROBLEMA
s=slide(WHITE)
box(s,0.7,0.5,12,0.9,"O problema atual",size=36,color=DEEP,bold=True,font=HSERIF)
items=[("1","Índice único","Reajuste só pelo CUB — índice de construção civil (R$/m²) aplicado a todas as modalidades."),
       ("2","Valor sem rastreabilidade","Não há modelo explícito de formação do valor-base por item."),
       ("3","Cobertura desatualizada","Faltam serviços novos (fotovoltaica) e unidades adequadas (kVA, kWp, ha, m³).")]
y=1.9
for n,t,d in items:
    circle_icon(s,0.8,y,0.7,TEAL,n)
    box(s,1.8,y-0.05,10.6,0.5,t,size=20,color=INK,bold=True)
    box(s,1.8,y+0.5,10.6,0.8,d,size=15,color=MUT)
    y+=1.55
box(s,0.8,6.85,12,0.5,"O próprio projeto interno reconhece que o CUB “não acompanha mais a evolução dos preços”.",
    size=13,color=DEEP,italic=True)

# 3 OPORTUNIDADE / EVIDENCIA
s=slide(LIGHT)
box(s,0.7,0.5,12,0.9,"A oportunidade: evidência já disponível",size=34,color=DEEP,bold=True,font=HSERIF)
box(s,0.72,1.45,12,0.5,"Base de ARTs do CREA-BA (2022), agregada e anonimizada",size=15,color=MUT)
for i,(num,lbl) in enumerate([("726.028","linhas de atividade (ART)"),("230.928","ARTs distintas"),
        ("R$ 1.570","mediana do valor declarado"),("1.047","faixas Atividade × Unidade")]):
    x=0.8+i*3.05
    shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,2.35,2.85,2.0,fill=WHITE)
    stat(s,x+0.3,2.65,2.3,num,lbl)
box(s,0.8,4.9,12,0.5,"Faixa interquartil (P25–P75): R$ 540 a R$ 8.000 · 97% das ARTs na Bahia",size=15,color=INK,bold=True)
box(s,0.8,5.5,12,1.2,"111 unidades de medida — predominam m² e “unidade”, com forte volume de kWp (≈40 mil) e kVA "
    "(≈11 mil), confirmando a necessidade de unidades específicas por serviço.",size=14,color=MUT)

# 4 LOCALIZACAO
s=slide(WHITE)
box(s,0.7,0.5,12,0.9,"Localização — onde está a demanda",size=34,color=DEEP,bold=True,font=HSERIF)
box(s,0.72,1.45,12,0.5,"Concentração geográfica das ARTs (2022) — base do ajuste regional",size=15,color=MUT)
cidades=[("Salvador","85.334"),("Feira de Santana","26.121"),("Vitória da Conquista","20.982"),
         ("Camaçari","20.741"),("Barreiras","20.642"),("Luís Eduardo Magalhães","18.446")]
for i,(c,n) in enumerate(cidades):
    x=0.8+(i%3)*4.05; y=2.3+(i//3)*1.6
    shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,y,3.8,1.35,fill=LIGHT)
    box(s,x+0.3,y+0.2,3.3,0.5,n,size=24,color=TEAL,bold=True,font=HSERIF)
    box(s,x+0.3,y+0.78,3.3,0.4,c,size=13,color=INK)
box(s,0.8,5.9,12,0.8,"Eixos: Região Metropolitana de Salvador, oeste agrícola (Barreiras/LEM/São Desidério) "
    "e sul/extremo-sul. O painel interativo traz o mapa completo da Bahia.",size=14,color=MUT)

# 5 LIMITES
s=slide(MID)
box(s,0.7,0.6,12,0.9,"Os limites dos dados (honestidade metodológica)",size=30,color=WHITE,bold=True,font=HSERIF)
pts=["A ART NÃO é honorário: valor_contrato pode ser valor de obra, de contrato ou declarado.",
     "Há outliers e erros (máximo observado > R$ 800 milhões) e 111 unidades distintas.",
     "Por isso usamos MEDIANA e intervalo interquartil — nunca a média.",
     "Uso apenas agregado e anonimizado; supressão de células com n<5; sem ranking individual (LGPD)."]
y=2.0
for p in pts:
    circle_icon(s,0.8,y,0.5,ACC,"✓",gsize=14)
    box(s,1.55,y+0.02,10.9,0.8,p,size=16,color=RGBColor(0xE8,0xEE,0xF5))
    y+=1.15
box(s,0.8,6.7,12,0.5,"A evidência é INDIRETA — serve para calibrar faixas, não para provar preço.",
    size=14,color=ACC,bold=True,italic=True)

# 6 PROPOSTA 3 CAMADAS
s=slide(WHITE)
box(s,0.7,0.5,12,0.9,"A proposta: 3 camadas + calibração",size=34,color=DEEP,bold=True,font=HSERIF)
cam=[("Catálogo de atividades","Nível + Atividade do sistema CREA — o que se precifica."),
     ("Valor-base por esforço","Tempo técnico × valor-hora de referência + custos (piso profissional)."),
     ("Ajuste regional","Pesos socioeconômicos por região da Bahia (modelo já iniciado).")]
for i,(t,d) in enumerate(cam):
    x=0.8+i*4.05
    shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,2.0,3.8,2.7,fill=LIGHT)
    circle_icon(s,x+0.3,2.3,0.7,DEEP,str(i+1))
    box(s,x+0.3,3.15,3.2,0.6,t,size=17,color=INK,bold=True)
    box(s,x+0.3,3.8,3.2,0.8,d,size=13,color=MUT)
shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,0.8,5.15,11.7,1.4,fill=DEEP)
box(s,1.1,5.35,11.1,0.5,"+ Calibração / validação contra mediana e IQR das ARTs",size=18,color=WHITE,bold=True)
box(s,1.1,5.95,11.1,0.5,"Saída em FAIXAS: piso técnico (P25) · referência (mediana) · teto orientativo (P75)",
    size=14,color=RGBColor(0xCA,0xDC,0xFC))

# 7 FAIXAS
s=slide(LIGHT)
box(s,0.7,0.5,12,0.9,"Faixas, não preço único",size=34,color=DEEP,bold=True,font=HSERIF)
box(s,0.72,1.45,12,0.6,"Mais defensável tecnicamente e juridicamente; absorve a dispersão real do mercado.",
    size=15,color=MUT)
faix=[("PISO TÉCNICO","P25","Limite anti-aviltamento",TEAL),
      ("REFERÊNCIA","Mediana","Estimativa central",DEEP),
      ("TETO ORIENTATIVO","P75","Maior complexidade",MID)]
for i,(t,p,d,col) in enumerate(faix):
    x=0.8+i*4.05
    shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,2.4,3.8,2.4,fill=WHITE)
    shape(s,MSO_SHAPE.OVAL,x+1.5,2.7,0.8,0.8,fill=col)
    box(s,x+0.2,3.7,3.4,0.5,t,size=16,color=INK,bold=True,align=PP_ALIGN.CENTER)
    box(s,x+0.2,4.2,3.4,0.4,p+" — "+d,size=12.5,color=MUT,align=PP_ALIGN.CENTER)
box(s,0.8,5.4,12,0.9,"Exemplos calibrados (ART 2022): Projeto/m² → mediana R$ 1.500 (IQR 800–5.000) · "
    "Projeto/kWp → R$ 1.000 (552–2.000).",size=14,color=INK)

# 8 ENTREGAVEIS
s=slide(WHITE)
box(s,0.7,0.5,12,0.9,"O que já está pronto",size=34,color=DEEP,bold=True,font=HSERIF)
ent=[("Painel interativo","Dashboard HTML com mapa da Bahia e gráficos — abre em 1 clique."),
     ("Planilha-modelo","1.047 itens (Atividade × Unidade) com faixas observadas."),
     ("Diagnósticos","Metodologia atual e dados de ART, com lacunas e riscos."),
     ("Documentos institucionais","Proposta, nota técnica, minuta de recomendação e roteiro.")]
for i,(t,d) in enumerate(ent):
    x=0.8+(i%2)*6.05; y=1.95+(i//2)*2.2
    shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,y,5.8,1.9,fill=LIGHT)
    circle_icon(s,x+0.35,y+0.35,0.7,TEAL,"●",gsize=14)
    box(s,x+1.3,y+0.35,4.3,0.5,t,size=18,color=INK,bold=True)
    box(s,x+1.3,y+0.9,4.3,0.85,d,size=13,color=MUT)

# 9 GOVERNANCA
s=slide(WHITE)
box(s,0.7,0.5,12,0.9,"Governança, LGPD e segurança jurídica",size=32,color=DEEP,bold=True,font=HSERIF)
g=[("Revisão anual","Comissão presidida pelo SENGE; manter / atualizar / inserir itens."),
   ("Transparência","Nota técnica, versionamento e changelog a cada revisão."),
   ("LGPD","Apenas agregados; supressão n<5; sem ranking nem dados pessoais."),
   ("Caráter orientativo","Referência técnica anti-aviltamento — não tabelamento.")]
y=1.95
for t,d in g:
    circle_icon(s,0.8,y,0.6,DEEP,"›")
    box(s,1.65,y,5.0,0.5,t,size=18,color=INK,bold=True)
    box(s,6.7,y+0.02,5.8,0.9,d,size=14,color=MUT)
    y+=1.25

# 10 PROXIMOS PASSOS
s=slide(LIGHT)
box(s,0.7,0.5,12,0.9,"Próximos passos",size=34,color=DEEP,bold=True,font=HSERIF)
steps=[("1","Aprovar a arquitetura metodológica (3 camadas + calibração)."),
       ("2","Constituir a comissão e definir parâmetros (valor-hora, tempo técnico)."),
       ("3","Pesquisa de preços com entidades (≥5 por item)."),
       ("4","Validação jurídica do caráter orientativo e das normas."),
       ("5","Publicar a v1 com painel de evidências e planilha-modelo.")]
y=1.95
for n,d in steps:
    circle_icon(s,0.85,y,0.6,TEAL,n)
    box(s,1.7,y+0.05,10.8,0.6,d,size=16,color=INK)
    y+=1.0

# 11 ENCERRAMENTO
s=slide(MID)
shape(s,MSO_SHAPE.OVAL,-1.2,5.6,3.4,3.4,fill=DEEP)
box(s,0.9,2.5,11.5,1.6,"Modernizar a tabela sem inventar números:\ncada valor terá origem rastreável "
    "e cada faixa, respaldo empírico.",size=28,color=WHITE,bold=True,font=HSERIF,line=1.1)
box(s,0.92,4.5,11.5,0.6,"Valorização da engenharia baiana com segurança técnica e jurídica.",
    size=17,color=RGBColor(0xCA,0xDC,0xFC))
box(s,0.92,6.6,11.5,0.5,"SENGE/BA · Proposta de subsídio técnico · 2026",size=12,color=RGBColor(0x9F,0xB0,0xC3))

outp=os.path.join(BASE,'APRESENTACAO_CREA_SENGE.pptx')
prs.save(outp); print('WROTE',outp,'| slides:',len(prs.slides._sldIdLst))

